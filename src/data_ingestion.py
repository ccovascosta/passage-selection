import os
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation

def extract_data_from_pdf(file_path):
    text = ""
    metadata = {
        'filename': os.path.basename(file_path),
        'type': 'pdf',
        'title': None,
        'authors': None,
        'lastmodifiedtime': None
    }
    try:
        reader = PdfReader(file_path)
        for page in reader.pages:
            text += page.extract_text()
        info = reader.metadata
        metadata['title'] = info.get('/Title', None)
        metadata['authors'] = info.get('/Author', None)
        metadata['lastmodifiedtime'] = info.get('/ModDate', None)
    except Exception as e:
        print(f"Error extracting text and metadata from PDF {file_path}: {e}")
    return text, metadata

def extract_data_from_docx(file_path):
    text = ""
    metadata = {
        'filename': os.path.basename(file_path),
        'type': 'docx',
        'title': None,
        'authors': None,
        'lastmodifiedtime': None
    }
    try:
        doc = Document(file_path)
        text = "\n".join([para.text for para in doc.paragraphs])
        core_properties = doc.core_properties
        metadata['title'] = core_properties.title
        metadata['authors'] = core_properties.author
        metadata['lastmodifiedtime'] = core_properties.modified
    except Exception as e:
        print(f"Error extracting text and metadata from DOCX {file_path}: {e}")
    return text, metadata

def extract_data_from_pptx(file_path):
    text = ""
    metadata = {
        'filename': os.path.basename(file_path),
        'type': 'pptx',
        'title': None,
        'authors': None,
        'lastmodifiedtime': None
    }
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        core_properties = prs.core_properties
        metadata['title'] = core_properties.title
        metadata['authors'] = core_properties.author
        metadata['lastmodifiedtime'] = core_properties.modified
    except Exception as e:
        print(f"Error extracting text and metadata from PPTX {file_path}: {e}")
    return text, metadata

def load_document(file_path):
    if file_path.endswith('.pdf'):
        return extract_data_from_pdf(file_path)
    elif file_path.endswith('.docx'):
        return extract_data_from_docx(file_path)
    elif file_path.endswith('.pptx'):
        return extract_data_from_pptx(file_path)
    else:
        print(f"Unsupported file format: {file_path}")
        return None, None
